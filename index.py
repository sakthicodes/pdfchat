import streamlit as st
from dotenv import load_dotenv
import os
import asyncio
from langchain_community.document_loaders import PyPDFLoader
from langchain_openai import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.chains import RetrievalQA
from langchain_community.chat_models import ChatOpenAI
from langchain_openai import ChatOpenAI
from docx import Document
from pptx import Presentation
from fpdf import FPDF
from openai import OpenAI
load_dotenv()

openai_api_key = os.getenv("openai_api_key")
client = OpenAI(api_key=openai_api_key)
# Initialize chat history
if "messages" not in st.session_state:
    st.session_state.messages = []

if "openai_model" not in st.session_state:
    st.session_state["openai_model"] = "gpt-3.5-turbo"
# Initialize new_db variable
new_db = None

# Function to process uploaded document asynchronously
async def process_uploaded_document_async(uploaded_file):
    global new_db  # Use the global new_db variable
    if uploaded_file is not None:
        # Save the uploaded file temporarily
        with open("temp_file", "wb") as f:
            f.write(uploaded_file.read())

        # Convert non-PDF files to PDF
        file_extension = uploaded_file.name.split(".")[-1].lower()
        pdf_file_path = convert_to_pdf("temp_file", file_extension)

        # Load the document based on its type
        loader = PyPDFLoader(pdf_file_path)
        content = loader.load_and_split()  # Returns a list of page content strings

        # Process the document and create the FAISS index
        embeddings = OpenAIEmbeddings(openai_api_key=openai_api_key)
        texts = [content] if isinstance(content, str) else content  # Convert content to list if it's a string
        db = FAISS.from_documents(texts, embeddings)  # Change to from_documents
        db.save_local("faiss_index")
        new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)

        # Remove the temporary file
        os.remove("temp_file")

# Function to convert non-PDF files to PDF
def convert_to_pdf(file_path, file_extension):
    pdf_file_path = None
    if file_extension == "docx":
        doc = Document(file_path)
        pdf = FPDF()
        for paragraph in doc.paragraphs:
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            pdf.cell(200, 10, txt=paragraph.text.encode('latin-1', 'replace').decode('latin-1'), ln=True, align="L")
        pdf_file_path = file_path + ".pdf"
        pdf.output(pdf_file_path)
    elif file_extension == "pptx":
        prs = Presentation(file_path)
        pdf = FPDF()
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    pdf.add_page()
                    pdf.set_font("Arial", size=12)
                    pdf.cell(200, 10, txt=shape.text.encode('latin-1', 'replace').decode('latin-1'), ln=True, align="L")
        pdf_file_path = file_path + ".pdf"
        pdf.output(pdf_file_path)
    elif file_extension == "txt":
        with open(file_path, "r", encoding="utf-8") as txt_file:
            content = txt_file.read()
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        pdf.multi_cell(0, 10, txt=content.encode('latin-1', 'replace').decode('latin-1'))
        pdf_file_path = file_path + ".pdf"
        pdf.output(pdf_file_path)
    elif file_extension == "pdf":
        pdf_file_path = file_path
    else:
        st.error("Unsupported file type. Please upload a PDF, DOCX, PPTX, or TXT file.")
    return pdf_file_path

# Function to query the QA chain
def ask(user_query):
    global new_db
    if new_db is None:
        st.error("Please upload a document first to initialize the database.")
        return
    llm = ChatOpenAI(openai_api_key=openai_api_key)
    qa_chain = RetrievalQA.from_chain_type(llm, retriever=new_db.as_retriever())
    res = qa_chain({"query": user_query})
    result = res["result"]
    # Ensure result is always a list of dictionaries
    if isinstance(result, str):
        result = [{"role": "assistant", "content": result}]
    elif isinstance(result, dict):
        result = [result]
    return result

# Sidebar document uploader
uploaded_file = st.sidebar.file_uploader("Upload Document", type=["pdf", "docx", "pptx", "txt"])

# Process uploaded document and update database asynchronously
if uploaded_file:
    st.spinner("Processing document...")
    asyncio.run(process_uploaded_document_async(uploaded_file))
    st.success("Document processed successfully.")

# Display chat messages from history on app rerun
for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])
        
# React to user input
if prompt := st.chat_input("What is up?", key="user_input"):
    # Display user message in chat message container
    st.chat_message("user").markdown(prompt)
    # Add user message to chat history
    st.session_state.messages.append({"role": "user", "content": prompt})

    # Query the QA chain
    response1 = ask(prompt)
    with st.spinner("Thinking..."):
        # Display responses as they arrive
        for resp in response1:
            with st.chat_message("assistant"):
                st.markdown(resp["content"])
                st.session_state.messages.append({"role": "assistant", "content": resp["content"]})